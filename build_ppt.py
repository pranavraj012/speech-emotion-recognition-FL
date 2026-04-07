from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).parent
ASSETS = ROOT / "presentation_assets"
OUT = ROOT / "SER_FL_Presentation.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet


def add_image_slide(prs, title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.2), width=Inches(11.5))
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4))
        tx.text_frame.text = caption


def main():
    prs = Presentation()

    add_title_slide(
        prs,
        "Privacy-Preserving Speech Emotion Recognition",
        "Federated Learning simulation with Flower on RAVDESS",
    )

    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "Classify speech emotion while preserving user privacy",
            "Raw audio stays local on each client",
            "Target metric: at least 85% test accuracy",
        ],
    )

    add_bullet_slide(
        prs,
        "Data and Transformation",
        [
            "Dataset: RAVDESS speech clips",
            "6 target emotions: neutral, happy, sad, angry, fearful, disgusted",
            "HuBERT-Large generates 1024-dim embeddings per audio sample",
        ],
    )

    add_bullet_slide(
        prs,
        "Federated Environment",
        [
            "Framework: Flower",
            "Strategy: FedAvg",
            "Simulation: 5 local clients and 1 server",
            "Split policy: global 80/20 holdout, then train partitioning across clients",
        ],
    )

    add_bullet_slide(
        prs,
        "Model and Stack",
        [
            "Model: MLP (1024 -> 256 -> 6)",
            "Training: Adam optimizer, mini-batch local updates",
            "Stack: PyTorch, TorchAudio, Transformers, Flower, scikit-learn, Streamlit",
        ],
    )

    for title, filename in [
        ("Global Accuracy by Round", "accuracy_curve.png"),
        ("Global Loss by Round", "loss_curve.png"),
        ("Confusion Matrix", "confusion_matrix.png"),
        ("Per-Class Accuracy", "class_accuracy.png"),
    ]:
        image = ASSETS / filename
        if image.exists():
            add_image_slide(prs, title, image)

    add_bullet_slide(
        prs,
        "What I Learned",
        [
            "Federated evaluation requires strict holdout discipline",
            "Model choice strongly affects convergence in FL",
            "Privacy settings create a measurable utility tradeoff",
            "Good logging and visualization are essential for debugging FL",
        ],
    )

    add_bullet_slide(
        prs,
        "Demo and Future Work",
        [
            "Demo app: upload audio or record live voice in Streamlit",
            "Current run reached around 93% accuracy",
            "Future: stronger DP with minimal accuracy loss, true multi-machine FL",
        ],
    )

    prs.save(OUT)
    print(f"Saved PPT: {OUT}")


if __name__ == "__main__":
    main()
